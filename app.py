"""
RobChat: A RAG-based Document Question Answering System

This application implements a Retrieval-Augmented Generation (RAG) system that allows users to:
1. Upload various document types (PDF, DOCX, TXT, HTML, PPTX, XLSX)
2. Process and index document content
3. Query the documents using natural language
4. Get AI-generated answers with source citations

The system uses:
- LangChain for the RAG pipeline
- Llama2 (7B-Chat) for text generation
- sentence-transformers for embeddings
- ChromaDB for vector storage
- FastAPI for the web API
"""

from fastapi import FastAPI, HTTPException, UploadFile, File, Depends, status, Form
from fastapi.middleware.cors import CORSMiddleware
from fastapi.staticfiles import StaticFiles
from fastapi.responses import FileResponse
from fastapi.security import OAuth2PasswordBearer, OAuth2PasswordRequestForm
from pydantic import BaseModel
from typing import List, Optional, Dict, Any, Annotated
import os
import shutil
from pathlib import Path
import json
import logging
import time
import signal
import sys
from functools import wraps
from tenacity import retry, stop_after_attempt, wait_exponential
from langchain_community.document_loaders import DirectoryLoader, TextLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_community.vectorstores import Chroma
from langchain_community.llms import LlamaCpp
from langchain.chains import RetrievalQA, ConversationalRetrievalChain, LLMChain
from langchain.chains.combine_documents import create_stuff_documents_chain
from langchain.prompts import PromptTemplate, ChatPromptTemplate, HumanMessagePromptTemplate, SystemMessagePromptTemplate
import torch
from huggingface_hub import scan_cache_dir, hf_hub_download
from huggingface_hub.constants import HUGGINGFACE_HUB_CACHE
from docx import Document
from bs4 import BeautifulSoup
from pptx import Presentation
from openpyxl import load_workbook
import io
from langchain.memory import ConversationBufferMemory
import yaml
import pypdf
import traceback
from datetime import datetime, timedelta
from jose import JWTError, jwt
import bcrypt
from passlib.context import CryptContext
import chromadb
from chromadb.config import Settings
import threading

# -------------------------------------------------------------------------------------------------
# CONFIGURATION
# -------------------------------------------------------------------------------------------------

# Load configuration first, before setting up logging
def load_config() -> Dict[Any, Any]:
    """
    Load configuration from YAML file and check for required files.

    Returns:
        dict: Configuration dictionary

    Exits:
        If config.yaml or users.yaml doesn't exist
    """
    # Check for users.yaml first
    if not os.path.exists("users.yaml"):
        print("users.yaml file not found. Please create a users.yaml file based on users.example.yaml.")
        sys.exit(1)

    # Check for config.yaml q
    if not os.path.exists("config.yaml"):
        print("users.yaml file not found. Please create a users.yaml file based on users.example.yaml.")
        sys.exit(1)

    try:
        with open("config.yaml", 'r') as f:
            config = yaml.safe_load(f)
        return config
    except FileNotFoundError:
        print('Configuration file "config.yaml" not found. Please create a config.yaml file.')
        sys.exit(1)

# Load configuration before setting up logging
config = load_config()

# Setup logging with config
logging.basicConfig(
    level=getattr(logging, config.get("server", {}).get("log_level", "INFO").upper()),
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)
logger.info("Loaded configuration from config.yaml")

# Set specific loggers to INFO level
logging.getLogger('python_multipart.multipart').setLevel(logging.INFO)
logging.getLogger('urllib3.connectionpool').setLevel(logging.INFO)

# Device configuration
"""Configure the compute device (CPU, CUDA, or MPS) based on availability"""
try:
    if torch.backends.mps.is_available():
        DEVICE = "mps"  # Apple Silicon (M1/M2) GPU
        logger.info("Using Apple Metal (MPS) device")
    elif torch.cuda.is_available():
        DEVICE = "cuda"
        logger.info("Using CUDA device: %s", torch.cuda.get_device_name(0))
    else:
        DEVICE = "cpu"
        logger.info("Using CPU device")

    logger.info("PyTorch version: %s", torch.__version__)
    logger.info("Device: %s", DEVICE)
except Exception as e:
    logger.warning("Error detecting device, falling back to CPU: %s", str(e))
    DEVICE = "cpu"

# Update global constants from config
DATA_DIR = config["data"]["base_directory"]

# Define the prompts
HUMAN_TEMPLATE = "Question: {question}"

# API Models
class Query(BaseModel):
    text: str

class FileList(BaseModel):
    files: List[str]

class UserProjects(BaseModel):
    """Model for user projects response."""
    projects: List[str]  # List of all projects
    current_project: str

# Initialize FastAPI app
app = FastAPI(title="RobChat API")

# Setup static paths
static_path = os.path.join(os.path.dirname(__file__), "static")

# Mount the static files
if os.path.exists(static_path):
    app.mount("/static", StaticFiles(directory=static_path), name="static")

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Global caches
_llm_cache: Optional[LlamaCpp] = None
_vector_store_cache: Dict[str, Chroma] = {}  # Key format: "user:project"
_llm_lock = threading.Lock()  # Lock for thread-safe LLM initialization

def initialize_llm():
    """Initialize the LLM in a background thread"""
    global _llm_cache
    try:
        logger.info("Starting LLM initialization in background...")
        with _llm_lock:
            if _llm_cache is None:  # Double-check pattern
                _llm_cache = get_llm()
                logger.info("LLM initialized successfully")
            else:
                logger.info("LLM already initialized by another worker")
    except Exception as e:
        logger.error(f"Failed to initialize LLM: {str(e)}")
        # Don't raise here since we're in a thread

async def startup_event():
    """Start the LLM initialization in a background thread"""
    import threading
    thread = threading.Thread(target=initialize_llm)
    thread.daemon = True  # Make thread daemon so it doesn't prevent shutdown
    thread.start()
    logger.info("Started LLM initialization thread")

app.add_event_handler("startup", startup_event)

# -------------------------------------------------------------------------------------------------
# AUTHENTICATION
# -------------------------------------------------------------------------------------------------

# Authentication models
class Token(BaseModel):
    access_token: str
    token_type: str

class TokenData(BaseModel):
    username: Optional[str] = None

class User(BaseModel):
    username: str
    email: str
    fullname: str
    admin: bool
    disabled: Optional[bool] = None

class UserInDB(User):
    hashed_password: str

# Authentication settings
SECRET_KEY = config["auth"]["secret_key"]
ALGORITHM = config["auth"]["algorithm"]
ACCESS_TOKEN_EXPIRE_MINUTES = config["auth"]["access_token_expire_minutes"]

pwd_context = CryptContext(schemes=["bcrypt"], deprecated="auto")
oauth2_scheme = OAuth2PasswordBearer(tokenUrl="token")

def get_password_hash(password):
    return bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')

def verify_password(plain_password, hashed_password):
    if isinstance(hashed_password, str):
        hashed_password = hashed_password.encode('utf-8')
    return bcrypt.checkpw(plain_password.encode('utf-8'), hashed_password)

def load_users():
    """Load users from users.yaml file"""
    try:
        with open('users.yaml', 'r') as f:
            users = yaml.safe_load(f)
            return users
    except FileNotFoundError:
        logger.error("users.yaml file not found")
        return {}
    except Exception as e:
        logger.error("Error loading users.yaml: %s", str(e))
        return {}

def get_user(username: str):
    """Get user from users.yaml file"""
    users = load_users()
    if username in users:
        user_dict = users[username]
        user_dict["username"] = username
        user_dict["hashed_password"] = user_dict.pop("password")  # Map password to hashed_password
        return UserInDB(**user_dict)
    return None

def authenticate_user(username: str, password: str):
    """Authenticate user against users.yaml file"""
    user = get_user(username)
    if not user:
        return False
    if not verify_password(password, user.hashed_password):
        return False
    return user

def create_access_token(data: dict, expires_delta: Optional[timedelta] = None):
    to_encode = data.copy()
    if expires_delta:
        expire = datetime.utcnow() + expires_delta
    else:
        expire = datetime.utcnow() + timedelta(minutes=15)
    to_encode.update({"exp": expire})
    encoded_jwt = jwt.encode(to_encode, SECRET_KEY, algorithm=ALGORITHM)
    return encoded_jwt

async def get_current_user(token: Annotated[str, Depends(oauth2_scheme)]):
    credentials_exception = HTTPException(
        status_code=status.HTTP_401_UNAUTHORIZED,
        detail="Could not validate credentials",
        headers={"WWW-Authenticate": "Bearer"},
    )
    try:
        payload = jwt.decode(token, SECRET_KEY, algorithms=[ALGORITHM])
        username: str = payload.get("sub")
        if username is None:
            raise credentials_exception
        token_data = TokenData(username=username)
    except JWTError as exc:
        raise credentials_exception from exc
    user = get_user(username=token_data.username)
    if user is None:
        raise credentials_exception
    return user

async def get_current_active_user(
    current_user: Annotated[User, Depends(get_current_user)]
):
    if current_user.disabled:
        raise HTTPException(status_code=400, detail="Inactive user")
    return current_user

# -------------------------------------------------------------------------------------------------
# RAG COMPONENTS
# -------------------------------------------------------------------------------------------------

# Retry decorator for model operations
def retry_with_fallback(max_attempts=3):
    """
    Decorator that implements retry logic with exponential backoff.

    Args:
        max_attempts (int): Maximum number of retry attempts

    Returns:
        Function wrapper that implements the retry logic
    """
    def decorator(func):
        @wraps(func)
        def wrapper(*args, **kwargs):
            for attempt in range(max_attempts):
                try:
                    return func(*args, **kwargs)
                except Exception as e:
                    logger.warning(f"Attempt {attempt + 1}/{max_attempts} failed: {str(e)}")
                    if attempt == max_attempts - 1:  # Last attempt
                        logger.error(f"All attempts failed for {func.__name__}: {str(e)}")
                        raise
                    time.sleep(2 ** attempt)  # Exponential backoff
            return None
        return wrapper
    return decorator

# Initialize the RAG components
@retry_with_fallback(max_attempts=3)
def get_embeddings():
    """
    Initialize and return the sentence transformer embeddings model.
    Uses model specified in configuration.
    """
    model_name = config["models"]["embeddings"]["model_name"]
    model_kwargs = {'device': DEVICE}
    embeddings = HuggingFaceEmbeddings(model_name=model_name, model_kwargs=model_kwargs)
    return embeddings

@retry_with_fallback(max_attempts=3)
def get_llm():
    """
    Initialize and return the Llama2 language model.
    Uses model and parameters specified in configuration.
    Returns cached instance if available.
    """
    global _llm_cache
    if _llm_cache is not None:
        return _llm_cache

    try:
        llm_config = config["models"]["llm"]
        logger.info("Downloading model from Hugging Face...")
        model_path = hf_hub_download(
            repo_id=llm_config["repo_id"],
            filename=llm_config["filename"],
            local_files_only=False,  # Ensure we download if not present
            force_download=False,     # Don't re-download if already present
            resume_download=True      # Resume interrupted downloads
        )
        logger.info("Model downloaded, initializing...")

        n_gpu_layers = 1 if DEVICE == "mps" else 0

        _llm_cache = LlamaCpp(
            model_path=model_path,
            temperature=llm_config["temperature"],
            max_tokens=llm_config["max_tokens"],
            n_ctx=llm_config["context_window"],
            top_p=llm_config["top_p"],
            n_gpu_layers=n_gpu_layers,
            n_batch=llm_config["n_batch"],
            verbose=True,
        )

        # Test the model to ensure it's fully loaded
        logger.info("Testing model initialization...")
        test_prompt = "Hello, are you ready?"
        _llm_cache(test_prompt)
        logger.info("Model initialization complete")

        return _llm_cache
    except Exception as e:
        logger.error(f"Error initializing LLM: {str(e)}")
        raise

# Initialize vector store
@retry_with_fallback(max_attempts=3)
def get_vector_store(user: str, project: str, load_documents: bool = False):
    """
    Initialize or load the vector store for a user/project.
    Returns cached instance if available for the same user/project.

    Args:
        user (str): User identifier
        project (str): Project identifier
        load_documents (bool): Whether to load documents from disk (default: False)

    Returns:
        Chroma: The vector store instance
    """
    global _vector_store_cache
    cache_key = f"{user}:{project}"

    # Check if we have a cached instance for this user/project
    if cache_key in _vector_store_cache:
        vector_store = _vector_store_cache[cache_key]
        # If load_documents is True and store is empty, we need to reload
        if not (load_documents and vector_store._collection.count() == 0):
            return vector_store

    try:
        project_path = os.path.join(DATA_DIR, user, project)
        if not os.path.exists(project_path):
            os.makedirs(project_path)

        persist_directory = os.path.join(project_path, "chroma_db")
        embeddings = get_embeddings()

        # Initialize ChromaDB client
        client = chromadb.PersistentClient(path=persist_directory)

        # Get or create collection with fixed name "robchat"
        try:
            collection = client.get_collection(name="robchat")
        except:
            collection = client.create_collection(name="robchat")

        # Initialize vector store with the collection
        vector_store = Chroma(
            persist_directory=persist_directory,
            embedding_function=embeddings,
            collection_name="robchat"
        )

        # Only load documents if explicitly requested and the store is empty
        if load_documents and vector_store._collection.count() == 0:
            logger.info(f"Loading documents from {project_path}")
            # Load all text files from the project directory
            loader = DirectoryLoader(
                project_path,
                glob="**/*.txt",
                loader_cls=TextLoader
            )
            documents = loader.load()

            if documents:
                logger.info(f"Found {len(documents)} documents to process")
                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=1000,
                    chunk_overlap=200
                )
                texts = text_splitter.split_documents(documents)
                vector_store.add_documents(texts)
                logger.info(f"Added {len(texts)} text chunks to vector store")
            else:
                logger.info("No documents found in project directory")
        else:
            logger.info(f"Using existing vector store with {vector_store._collection.count()} documents")

        # Cache the vector store
        _vector_store_cache[cache_key] = vector_store
        return vector_store

    except Exception as e:
        logger.error(f"Error initializing vector store: {str(e)}")
        raise

# -------------------------------------------------------------------------------------------------
# WEB PAGES
# -------------------------------------------------------------------------------------------------

@app.get("/")
async def read_root():
    """Serve the main application page"""
    if _llm_cache is None:
        return FileResponse("static/startup.html")
    return FileResponse("static/login.html")

@app.get("/chat")
async def read_chat():
    """Serve the chat interface page"""
    if _llm_cache is None:
        return FileResponse("static/startup.html")
    return FileResponse("static/chat.html")

# -------------------------------------------------------------------------------------------------
# AUTHENTICATION ENDPOINTS
# -------------------------------------------------------------------------------------------------

# API Endpoints
@app.post("/token", response_model=Token)
async def login_for_access_token(
    form_data: Annotated[OAuth2PasswordRequestForm, Depends()] = None,
    username: str = Form(None),
    password: str = Form(None)
):
    # Handle both OAuth2 form and regular form data
    if form_data:
        username = form_data.username
        password = form_data.password
    elif not username or not password:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Username and password are required"
        )

    user = authenticate_user(username, password)
    if not user:
        raise HTTPException(
            status_code=status.HTTP_401_UNAUTHORIZED,
            detail="Incorrect username or password",
            headers={"WWW-Authenticate": "Bearer"},
        )
    access_token_expires = timedelta(minutes=ACCESS_TOKEN_EXPIRE_MINUTES)
    access_token = create_access_token(
        data={
            "sub": user.username,
            "fullname": user.fullname
        },
        expires_delta=access_token_expires
    )
    return {"access_token": access_token, "token_type": "bearer"}

# -------------------------------------------------------------------------------------------------
# PROJECT ENDPOINTS
# -------------------------------------------------------------------------------------------------

@app.get("/api/{user}/projects")
async def get_user_projects(
    user: str,
    current_user: Annotated[User, Depends(get_current_active_user)]
):
    """
    Get all projects for a user and their current selected project.

    Args:
        user (str): User identifier

    Returns:
        UserProjects: List of all projects and current project
    """
    # Check if the authenticated user matches the requested user
    if current_user.username != user:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Permission denied: Cannot access projects of another user"
        )

    try:
        user_dir = os.path.join(DATA_DIR, user)
        if not os.path.exists(user_dir):
            return UserProjects(projects=[], current_project="default")

        # Get all project directories
        projects = [d for d in os.listdir(user_dir) if os.path.isdir(os.path.join(user_dir, d))]

        if not projects:
            return UserProjects(projects=[], current_project="default")

        # Get current project from the most recently accessed project directory
        current_project = ""
        latest_time = 0
        for project in projects:
            project_path = os.path.join(user_dir, project)
            access_time = os.path.getatime(project_path)
            if access_time > latest_time:
                latest_time = access_time
                current_project = project

        return UserProjects(projects=projects, current_project=current_project)
    except Exception as e:
        logger.error("Error getting user projects: %s", str(e))
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to get user projects: {str(e)}"
        ) from e

# -------------------------------------------------------------------------------------------------
# FILE HANDLING ENDPOINTS
# -------------------------------------------------------------------------------------------------

def parse_pdf(file_path: str) -> str:
    """
    Extract text content from PDF files with robust error handling.

    Args:
        file_path (str): Path to the PDF file

    Returns:
        str: Extracted text content

    Raises:
        ValueError: If PDF is encrypted or no text could be extracted
    """
    text = []
    try:
        with open(file_path, 'rb') as file:
            # Create a PDF reader object
            reader = pypdf.PdfReader(file)

            # Check if the PDF is encrypted
            if reader.is_encrypted:
                try:
                    reader.decrypt('')  # Try empty password
                except:
                    raise ValueError("PDF is encrypted and cannot be read")

            # Process each page
            for page_num in range(len(reader.pages)):
                try:
                    page = reader.pages[page_num]
                    page_text = page.extract_text()
                    if page_text.strip():  # Only add non-empty pages
                        text.append(page_text)
                except Exception as e:
                    logging.warning(f"Error extracting text from page {page_num}: {str(e)}")
                    text.append(f"[Error reading page {page_num}]")
                    continue

            if not text:  # If no text was extracted
                raise ValueError("No readable text found in PDF")

            return '\n'.join(text)
    except Exception as e:
        logging.error(f"Error processing PDF {file_path}: {str(e)}")
        raise ValueError(f"Failed to process PDF: {str(e)}")

def process_file_for_vectorstore(file_path: str) -> List[str]:
    """
    Process a file and prepare it for the vector store with proper chunking.

    Args:
        file_path (str): Path to the file to process

    Returns:
        List[str]: List of text chunks ready for vectorization

    Raises:
        ValueError: If file type is unsupported or processing fails
    """
    try:
        # Get file extension
        file_extension = os.path.splitext(file_path)[1].lower()

        # Extract content based on file type
        content = ""
        if file_extension == '.pdf':
            content = parse_pdf(file_path)
        elif file_extension == '.docx':
            doc = Document(file_path)
            content = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
        elif file_extension == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
        elif file_extension == '.html':
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f.read(), 'html.parser')
                content = soup.get_text(separator='\n')
        elif file_extension == '.pptx':
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            content = '\n'.join(texts)
        elif file_extension == '.xlsx':
            wb = load_workbook(file_path)
            texts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    row_text = ' '.join(str(cell) for cell in row if cell is not None)
                    if row_text.strip():
                        texts.append(row_text)
            content = '\n'.join(texts)
        else:
            raise ValueError(f"Unsupported file type: {file_extension}")

        if not content.strip():
            raise ValueError("No text content found in document")

        # Use chunk settings from config
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=config["vector_store"]["chunk_size"],
            chunk_overlap=config["vector_store"]["chunk_overlap"],
            length_function=len,
            separators=["\n\n", "\n", ".", " ", ""]
        )

        # Split text into chunks
        texts = text_splitter.split_text(content)

        if not texts:
            raise ValueError("No text chunks were generated from the document")

        return texts

    except Exception as e:
        logger.error(f"Error processing file for vector store: {str(e)}")
        raise

@app.get("/api/{user}/{project}/files")
async def list_files(
    user: str,
    project: str,
    current_user: Annotated[User, Depends(get_current_active_user)]
):
    """List all files in the project directory."""
    # Check if the authenticated user matches the requested user
    if current_user.username != user:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Permission denied: Cannot access files of another user"
        )

    try:
        project_path = os.path.join(DATA_DIR, user, project)
        if not os.path.exists(project_path):
            return {"files": []}

        files = []
        supported_extensions = config["data"]["supported_extensions"]
        for filename in os.listdir(project_path):
            if filename.endswith(tuple(supported_extensions)):
                file_path = os.path.join(project_path, filename)
                files.append({
                    "filename": filename,
                    "size": os.path.getsize(file_path),
                    "modified": datetime.fromtimestamp(os.path.getmtime(file_path)).isoformat()
                })

        return {"files": files}

    except Exception as e:
        logger.error("Error listing files: %s", str(e))
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to list files: {str(e)}"
        ) from e

@app.post("/api/{user}/{project}/files")
async def create_file(
    user: str,
    project: str,
    current_user: Annotated[User, Depends(get_current_active_user)],
    file: UploadFile = File(...)
):
    """
    Create (upload) a new file in the project.

    Args:
        user (str): User identifier
        project (str): Project identifier
        file (UploadFile): The file to upload

    Returns:
        dict: Information about the processed file

    Raises:
        HTTPException: For unsupported file types or processing errors
    """
    # Check if the authenticated user matches the requested user
    if current_user.username != user:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Permission denied: Cannot upload files for another user"
        )

    file_path = None
    replaced_existing = False
    try:
        # Get supported extensions from config
        supported_extensions = config["data"]["supported_extensions"]
        file_extension = os.path.splitext(file.filename)[1].lower()
        if file_extension not in supported_extensions:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"Unsupported file type. Supported types are: {', '.join(supported_extensions)}"
            )

        # Create project directory if it doesn't exist
        project_path = os.path.join(DATA_DIR, user, project)
        os.makedirs(project_path, exist_ok=True)

        # Check if file already exists and get its metadata
        file_path = os.path.join(project_path, file.filename)
        vector_store = get_vector_store(user, project)

        # Get all existing documents with this source file
        if vector_store._collection.count() > 0:
            existing_docs = vector_store._collection.get(
                where={"source": file_path}
            )
            if existing_docs and existing_docs['ids']:
                # Delete existing chunks for this file
                logger.info(f"Found existing chunks for {file.filename}, removing them...")
                vector_store._collection.delete(
                    ids=existing_docs['ids']
                )
                replaced_existing = True
                logger.info(f"Removed {len(existing_docs['ids'])} existing chunks")

        # Save the new file
        with open(file_path, "wb") as buffer:
            content = await file.read()
            buffer.write(content)

        # Get initial document count
        initial_count = vector_store._collection.count()
        logger.info(f"Initial document count in vector store: {initial_count}")

        # Process the file for vector store
        texts = process_file_for_vectorstore(file_path)
        logger.info(f"Successfully processed {file.filename} into {len(texts)} chunks")

        # Create documents from the text chunks
        documents = []
        for chunk in texts:
            documents.append({
                "page_content": chunk,
                "metadata": {"source": file_path}
            })

        # Add documents to vector store
        vector_store.add_texts(
            texts=[doc["page_content"] for doc in documents],
            metadatas=[doc["metadata"] for doc in documents]
        )

        # Get final document count
        final_count = vector_store._collection.count()
        logger.info(f"Added {len(texts)} chunks to vector store. Document count: {initial_count} -> {final_count}")

        return {
            "filename": file.filename,
            "size": os.path.getsize(file_path),
            "chunks": len(texts),
            "initial_count": initial_count,
            "final_count": final_count,
            "replaced_existing": replaced_existing
        }

    except HTTPException as he:
        raise he
    except Exception as e:
        # Clean up on any other error
        if file_path and os.path.exists(file_path):
            os.remove(file_path)
        logger.exception("Unexpected error during file upload: %s", str(e))
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"Failed to upload file: {str(e)}"
        ) from e


@app.delete("/api/{user}/{project}/files/{filename}")
async def delete_file(
    user: str,
    project: str,
    filename: str,
    current_user: Annotated[User, Depends(get_current_active_user)]
):
    """
    Delete a file and its content from the vector store.

    Args:
        user (str): User identifier
        project (str): Project identifier
        filename (str): Name of the file to delete

    Returns:
        dict: Status message
    """
    # Check if the authenticated user matches the requested user
    if current_user.username != user:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Permission denied: Cannot delete files of another user"
        )

    try:
        # Get the project path
        project_path = os.path.join(DATA_DIR, user, project)
        if not os.path.exists(project_path):
            raise HTTPException(status_code=404, detail="Project not found")

        # Get the file path
        file_path = os.path.join(project_path, filename)
        if not os.path.exists(file_path):
            raise HTTPException(status_code=404, detail="File not found")

        # Get the vector store
        vector_store = get_vector_store(user, project)

        # Delete the file from disk
        os.remove(file_path)

        # Delete the file's content from ChromaDB
        # We need to find all documents that came from this file
        # and delete them from the vector store
        collection = vector_store._collection
        results = collection.get(where={"source": file_path})

        if results and results["ids"]:
            collection.delete(ids=results["ids"])

        return {"status": "success", "message": f"File {filename} deleted successfully"}

    except HTTPException:
        raise
    except Exception as e:
        logger.error("Error deleting file: %s", str(e))
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=str(e)
        ) from e

# -------------------------------------------------------------------------------------------------
# QUERY ENDPOINTS
# -------------------------------------------------------------------------------------------------

@app.post("/api/{user}/{project}/query")
async def query_rag(
    query: Query,
    user: str,
    project: str,
    current_user: Annotated[User, Depends(get_current_active_user)]
):
    """
    Process a natural language query using RAG.

    Args:
        query (Query): The query to process
        user (str): User identifier
        project (str): Project identifier

    Returns:
        dict: The answer and sources
    """
    # Check if the authenticated user matches the requested user
    if current_user.username != user:
        raise HTTPException(
            status_code=status.HTTP_403_FORBIDDEN,
            detail="Permission denied: Cannot query documents of another user"
        )

    try:
        if not query.text.strip():
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="Query cannot be empty"
            )

        start_time = time.time()

        # Get vector store for this user/project
        vector_store = get_vector_store(user, project, load_documents=True)
        if vector_store._collection.count() == 0:
            raise HTTPException(status_code=404, detail="No documents found in vector store")

        # Create retriever with default k=4 if not specified in config
        num_sources = config.get("vector_store", {}).get("retrieval_k", 4)
        retriever = vector_store.as_retriever(
            search_kwargs={"k": num_sources}
        )
        logger.debug("Using retriever with k=%d sources", num_sources)

        # Get LLM from cache
        llm = get_llm()

        # Define custom prompt template
        template = config["rag"]["qa_template"]

        PROMPT = PromptTemplate(
            template=template,
            input_variables=["context", "question"]
        )

        qa_chain = RetrievalQA.from_chain_type(
            llm=llm,
            chain_type="stuff",
            retriever=retriever,
            chain_type_kwargs={"prompt": PROMPT},
            return_source_documents=True
        )

        # Execute query
        result = qa_chain({"query": query.text})

        # Calculate metrics
        end_time = time.time()
        elapsed_seconds = end_time - start_time

        # Get token counts from the LLM's last call
        try:
            input_tokens = llm.get_num_tokens(query.text)
            output_tokens = llm.get_num_tokens(result["result"])
            total_tokens = input_tokens + output_tokens
        except Exception as e:
            logger.warning(f"Failed to get token counts: {str(e)}")
            input_tokens = output_tokens = total_tokens = 0

        # Debug log the complete result with metrics
        logger.debug("Query: %s", query.text)
        logger.debug("Answer: %s", result["result"])
        logger.debug("Source documents: %s", [doc.metadata for doc in result["source_documents"]])
        logger.debug("Metrics - Time: %.2f seconds, Total Tokens: %d (Input: %d, Output: %d)",
                    elapsed_seconds, total_tokens, input_tokens, output_tokens)

        # Format response with deduplication
        seen_sources = set()  # Track unique sources
        sources = []
        source_number = 1

        for doc in result["source_documents"]:
            source_text = doc.page_content.strip()
            source_file = os.path.basename(doc.metadata.get('source', 'Unknown source'))

            # Create a unique key for this source using both content and filename
            source_key = f"{source_file}:{source_text}"

            # Only add if we haven't seen this exact source before
            if source_key not in seen_sources:
                seen_sources.add(source_key)
                sources.append({
                    "number": source_number,
                    "file": source_file,
                    "text": source_text,
                    "preview": source_text[:150] + "..." if len(source_text) > 150 else source_text
                })
                source_number += 1

        logger.info(f"Successfully processed query with {len(sources)} unique source documents in {elapsed_seconds:.2f} seconds")
        return {
            "answer": result["result"],
            "sources": sources,
            "metrics": {
                "time_seconds": round(elapsed_seconds, 2),
                "total_tokens": total_tokens,
                "input_tokens": input_tokens,
                "output_tokens": output_tokens
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error("Error processing query: %s", str(e))
        raise HTTPException(
            status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
            detail=f"An error occurred while processing your query: {str(e)}"
        ) from e


@app.get("/api/me", response_model=User)
async def get_current_user_details(
    current_user: Annotated[User, Depends(get_current_active_user)]
):
    """Get current user details"""
    return current_user

if __name__ == "__main__":
    import uvicorn
    import signal
    import sys

    # Configure uvicorn logging
    uvicorn_logger = logging.getLogger("uvicorn")
    uvicorn_logger.setLevel(logging.INFO)

    # Global flag to control server state
    should_exit = False

    def handle_exit(signum, frame):
        """Handle exit signals"""
        global should_exit
        should_exit = True
        logger.info("Shutdown requested. Stopping server...")

        # Clean up multiprocessing resources
        import multiprocessing
        multiprocessing.get_context()._cleanup()

        # Clean up LLM cache
        global _llm_cache
        if _llm_cache is not None:
            try:
                _llm_cache.__del__()
            except Exception as e:
                logger.warning(f"Error cleaning up LLM: {str(e)}")

        # Clean up vector store cache
        global _vector_store_cache
        _vector_store_cache.clear()

        sys.exit(0)

    # Register signal handlers
    signal.signal(signal.SIGINT, handle_exit)
    signal.signal(signal.SIGTERM, handle_exit)

    # Use server settings from config
    config = uvicorn.Config(
        app=app,
        host=config["server"]["host"],
        port=config["server"]["port"],
        log_level=config["server"]["log_level"]
    )
    server = uvicorn.Server(config)
    server.run()
